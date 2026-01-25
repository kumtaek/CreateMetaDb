"""
Markdown 형태의 발표자료를 PowerPoint(PPTX)로 변환하는 도구.

주의/한계
- Mermaid, 표, 코드블록 등은 PPT에서 '렌더링'하지 않고 텍스트로만 넣습니다.
- 슬라이드 구분은 `# 📄 Page N: ...` 헤더를 기준으로 합니다.

사용 예시
python tools/md_to_pptx.py --input "docs/AI코딩어시스턴스사례발표 - 발표자료.md" --output "temp/AI코딩어시스턴스사례발표.pptx"
"""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt


@dataclass
class SlideContent:
    """슬라이드 1장의 제목/본문 텍스트를 보관하는 모델."""

    title: str
    lines: list[str] = field(default_factory=list)


_PAGE_HEADER_RE = re.compile(r"^#\s*📄\s*Page\s*\d+\s*:\s*(.+?)\s*$")
_MD_HEADING_RE = re.compile(r"^(#{2,6})\s+(.*)$")
_MD_BULLET_RE = re.compile(r"^(\s*)[-*+]\s+(.*)$")


def _is_table_line(line: str) -> bool:
    """마크다운 테이블 라인 여부를 판정한다."""

    stripped = line.strip()
    return stripped.startswith("|") and stripped.endswith("|") and "|" in stripped[1:-1]


def _split_table_row(line: str) -> list[str]:
    """`| a | b |` 형태의 테이블 행을 셀 배열로 변환한다."""

    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _is_table_separator_row(cells: list[str]) -> bool:
    """`|---|---|` 같은 테이블 구분 행인지 판단한다."""

    if not cells:
        return False
    for cell in cells:
        normalized = cell.replace(":", "").replace("-", "").strip()
        if normalized != "":
            return False
    return True


def _flush_table_lines_as_text(table_lines: list[str], out_lines: list[str]) -> None:
    """누적된 테이블 라인을 사람이 읽기 쉬운 텍스트 라인으로 변환해 추가한다."""

    rows = [_split_table_row(line) for line in table_lines if _is_table_line(line)]
    if len(rows) < 2:
        out_lines.extend([line.rstrip() for line in table_lines if line.strip()])
        return

    header = rows[0]
    data_rows = []
    for row in rows[1:]:
        if _is_table_separator_row(row):
            continue
        data_rows.append(row)

    # 2열 테이블이면 '키: 값' 형태로, 그 외는 '헤더=값'들을 나열한다.
    for row in data_rows:
        if len(row) == 2:
            out_lines.append(f"• {row[0]}: {row[1]}")
            continue

        pairs: list[str] = []
        for idx, cell in enumerate(row):
            col = header[idx] if idx < len(header) else f"컬럼{idx + 1}"
            pairs.append(f"{col}={cell}")
        out_lines.append("• " + ", ".join(pairs))


def parse_markdown_pages(markdown_text: str) -> list[SlideContent]:
    """발표자료 마크다운을 'Page 헤더' 기준으로 슬라이드 목록으로 파싱한다."""

    slides: list[SlideContent] = []
    current: SlideContent | None = None

    in_code_block = False
    code_lang = ""
    code_lines: list[str] = []

    in_table = False
    table_lines: list[str] = []

    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip("\n")

        page_match = _PAGE_HEADER_RE.match(line)
        if page_match:
            # 슬라이드 경계 전에 테이블이 열려 있으면 flush
            if current and in_table and table_lines:
                _flush_table_lines_as_text(table_lines, current.lines)
                table_lines = []
                in_table = False

            title = page_match.group(0).lstrip("#").strip()
            current = SlideContent(title=title)
            slides.append(current)
            continue

        if current is None:
            continue

        if line.strip().startswith("```"):
            fence = line.strip()
            if not in_code_block:
                in_code_block = True
                code_lang = fence.strip("`").strip()
                code_lines = []
            else:
                in_code_block = False
                lang_suffix = f": {code_lang}" if code_lang else ""
                current.lines.append(f"[코드블록{lang_suffix}]")
                current.lines.extend([f"    {cl}" for cl in code_lines if cl.strip()])
                code_lang = ""
                code_lines = []
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        if _is_table_line(line):
            in_table = True
            table_lines.append(line)
            continue

        if in_table:
            # 테이블이 끝났으면 flush 후 현재 라인을 계속 처리
            if table_lines:
                _flush_table_lines_as_text(table_lines, current.lines)
            table_lines = []
            in_table = False

        stripped = line.strip()
        if not stripped or stripped == "---":
            continue

        if stripped.startswith(">"):
            quote = stripped.lstrip(">").strip()
            if quote:
                current.lines.append(f"“{quote}”")
            continue

        heading_match = _MD_HEADING_RE.match(line)
        if heading_match:
            heading_text = heading_match.group(2).strip()
            if heading_text:
                current.lines.append(heading_text)
            continue

        bullet_match = _MD_BULLET_RE.match(line)
        if bullet_match:
            indent = len(bullet_match.group(1).expandtabs(2))
            level = max(0, indent // 2)
            text = bullet_match.group(2).strip()
            if text:
                current.lines.append(f"{'  ' * level}• {text}")
            continue

        current.lines.append(stripped)

    if current and in_table and table_lines:
        _flush_table_lines_as_text(table_lines, current.lines)

    return slides


def _set_all_runs_font(text_frame, font_name: str, font_size_pt: int) -> None:
    """텍스트프레임 내 모든 런의 폰트를 지정한다."""

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)


def build_pptx(slides: Iterable[SlideContent], output_path: Path) -> None:
    """슬라이드 컨텐츠를 PPTX로 저장한다."""

    prs = Presentation()
    # 16:9 비율
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_only_layout = prs.slide_layouts[5]

    for slide_data in slides:
        slide = prs.slides.add_slide(title_only_layout)

        title_shape = slide.shapes.title
        title_shape.text = slide_data.title
        title_tf = title_shape.text_frame
        title_tf.paragraphs[0].font.size = Pt(34)
        title_tf.paragraphs[0].font.name = "Malgun Gothic"

        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.9), Inches(5.6))
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        # 빈 문단(1개)을 첫 줄로 사용
        if slide_data.lines:
            tf.paragraphs[0].text = slide_data.lines[0]
        else:
            tf.paragraphs[0].text = ""

        for line in slide_data.lines[1:]:
            tf.add_paragraph().text = line

        _set_all_runs_font(tf, font_name="Malgun Gothic", font_size_pt=18)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    """CLI 인자를 파싱한다."""

    parser = argparse.ArgumentParser(description="Markdown 발표자료를 PPTX로 변환합니다.")
    parser.add_argument(
        "--input",
        required=True,
        help="입력 Markdown 파일 경로",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="출력 PPTX 파일 경로",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    """엔트리포인트."""

    args = _parse_args(argv)
    input_path = Path(args.input)
    output_path = Path(args.output)

    if not input_path.exists():
        raise FileNotFoundError(f"입력 파일이 없습니다: {input_path}")

    markdown_text = input_path.read_text(encoding="utf-8")
    slides = parse_markdown_pages(markdown_text)

    if not slides:
        raise ValueError("슬라이드(Page 헤더)를 찾지 못했습니다. `# 📄 Page N: ...` 형식을 확인하세요.")

    build_pptx(slides, output_path=output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

